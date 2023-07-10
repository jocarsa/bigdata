from pptx import Presentation

ruta = "presentacion.pptx"
presentacion = Presentation(ruta)

contenido = []
for diapositiva in presentacion.slides:
    for forma in diapositiva.shapes:
        if hasattr(forma, "text"):
            contenido.append(forma.text)

print("\n".join(contenido))
